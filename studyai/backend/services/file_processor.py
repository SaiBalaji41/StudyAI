import io
import re
from pathlib import Path

import requests
from bs4 import BeautifulSoup
from werkzeug.datastructures import FileStorage

from config import ALLOWED_EXTENSIONS, MAX_FILE_SIZE


class FileProcessor:
    @staticmethod
    def allowed_file(filename: str) -> bool:
        return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

    @staticmethod
    def validate_file(file: FileStorage) -> tuple[bool, str]:
        if not file or not file.filename:
            return False, "No file provided"

        if not FileProcessor.allowed_file(file.filename):
            return False, f"Invalid file type. Allowed: {', '.join(ALLOWED_EXTENSIONS)}"

        file.seek(0, 2)
        size = file.tell()
        file.seek(0)

        if size > MAX_FILE_SIZE:
            return False, f"File too large. Maximum size: {MAX_FILE_SIZE // (1024 * 1024)} MB"

        if size == 0:
            return False, "File is empty"

        return True, ""

    @staticmethod
    def extract_text(file: FileStorage) -> tuple[str, str]:
        filename = file.filename or "unknown"
        ext = filename.rsplit(".", 1)[1].lower()
        file.seek(0)

        if ext == "txt":
            content = file.read().decode("utf-8", errors="ignore")
            return content.strip(), filename

        if ext == "pdf":
            return FileProcessor._extract_pdf(file), filename

        if ext == "docx":
            return FileProcessor._extract_docx(file), filename

        if ext == "pptx":
            return FileProcessor._extract_pptx(file), filename

        if ext in ("png", "jpg", "jpeg"):
            return FileProcessor._extract_image(file), filename

        raise ValueError(f"Unsupported file type: {ext}")

    @staticmethod
    def _extract_pdf(file: FileStorage) -> str:
        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(file.read()))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages).strip()

    @staticmethod
    def _extract_docx(file: FileStorage) -> str:
        from docx import Document

        doc = Document(io.BytesIO(file.read()))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs).strip()

    @staticmethod
    def _extract_pptx(file: FileStorage) -> str:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file.read()))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text.strip())
        return "\n\n".join(text_runs).strip()

    @staticmethod
    def _extract_image(file: FileStorage) -> str:
        file.seek(0)
        image_bytes = file.read()
        
        # Delayed import to avoid circular dependency
        from services.ai_service import ai_service
        filename = file.filename or "image.png"
        ext = filename.rsplit(".", 1)[1].lower()
        mime_type = f"image/{ext}" if ext != "jpg" else "image/jpeg"

        system_prompt = (
            "You are an expert OCR transcription engine. Extract all readable text "
            "from the provided image. Structure it logically. Do not add conversational intro/outro, "
            "return only the extracted text/markdown contents."
        )
        user_prompt = "Perform OCR on this image. Extract all text and diagrams description."

        try:
            extracted_text = ai_service._call_ai_multimodal(
                system_prompt, user_prompt, image_bytes, mime_type
            )
            return extracted_text.strip()
        except Exception as e:
            raise ValueError(f"Image analysis / OCR failed: {str(e)}")

    @staticmethod
    def extract_from_url(url: str) -> str:
        try:
            headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
            response = requests.get(url, headers=headers, timeout=15)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, "html.parser")
            
            for script in soup(["script", "style"]):
                script.decompose()
                
            content_tags = soup.find_all(["p", "h1", "h2", "h3", "h4", "li"])
            text = "\n\n".join([tag.get_text().strip() for tag in content_tags if tag.get_text().strip()])
            
            if not text:
                text = soup.get_text()
                
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            return "\n".join(chunk for chunk in chunks if chunk).strip()
        except Exception as e:
            raise ValueError(f"Failed to scrape URL: {str(e)}")

    @staticmethod
    def extract_from_youtube(url: str) -> str:
        from youtube_transcript_api import YouTubeTranscriptApi
        
        video_id_match = re.search(
            r'(?:v=|\/|embed\/|youtu\.be\/|v\/|watch\?v=|&v=)([^#\&\?]{11})', url
        )
        if not video_id_match:
            raise ValueError("Invalid YouTube URL. Could not parse video ID.")
            
        video_id = video_id_match.group(1)
        try:
            transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
            return " ".join([t["text"] for t in transcript_list]).strip()
        except Exception as e:
            raise ValueError(
                f"Could not retrieve transcript from YouTube video: {str(e)}. "
                "Make sure subtitles/transcripts are enabled for this video."
            )

    @staticmethod
    def validate_text_content(text: str) -> tuple[bool, str]:
        cleaned = text.strip()
        if not cleaned:
            return False, "Text content is empty"
        if len(cleaned) < 50:
            return False, "Text content is too short (minimum 50 characters)"
        if len(cleaned) > 50000:
            return False, "Text content is too long (maximum 50,000 characters)"
        return True, ""


file_processor = FileProcessor()
